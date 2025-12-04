"""PPTX to unified JSON extraction with Azure Document Intelligence for image OCR."""
from __future__ import annotations
import uuid
from typing import Dict, Any, List, Optional

from pptx import Presentation

from .helpers import (
    DIConfig,
    iter_text_shapes,
    iter_table_cells,
    iter_images,
    analyze_image_bytes,
)

# Re-export configs for convenience
__all__ = ["DIConfig", "pptx_to_unified_json"]


# ------------------------
# Public: main entry point
# ------------------------

def pptx_to_unified_json(
    pptx_path: str,
    di: DIConfig,
    *,
    include_tables: bool = True,
    compact: bool = True,
    verbose: bool = True,
) -> Dict[str, Any]:
    """
    Extract content from a PPTX file and return a unified JSON structure.
    
    Uses Azure Document Intelligence for image OCR.

    Args:
        pptx_path: Path to the PPTX file
        di: Document Intelligence configuration
        include_tables: Whether to include table cell text
        compact: If True, return simplified structure (recommended).
                 If False, include shape names and verbose OCR data.
        verbose: Print debug output
    
    Returns (compact=True):
    {
      "slides": [
        {
          "index": 1,
          "title": "Slide title if found",
          "text": ["Text block 1", "Text block 2", ...],
          "images": [
            {
              "name": "slide1_img1.png",
              "ocr_text": "Text extracted from image..."
            }
          ]
        }
      ],
      "unsupported_images": [...]
    }
    """
    if verbose:
        print(f"[DEBUG] Opening PPTX: {pptx_path}")
    prs = Presentation(pptx_path)
    if verbose:
        print(f"[DEBUG] PPTX opened, {len(prs.slides)} slides found")

    slides_out: List[Dict[str, Any]] = []
    unsupported_images: List[Dict[str, Any]] = []

    for i, slide in enumerate(prs.slides, start=1):
        if verbose:
            print(f"[DEBUG] Processing slide {i}...")
        
        if compact:
            slide_entry: Dict[str, Any] = {"index": i, "title": None, "text": [], "images": []}
        else:
            slide_entry: Dict[str, Any] = {"index": i, "text": [], "images": []}

        # 1) native text boxes
        text_shapes = list(iter_text_shapes(slide))
        for t in text_shapes:
            text_content = t.get("text", "").strip()
            if not text_content:
                continue
                
            if compact:
                # Try to identify title (first text shape or shape named "Title")
                shape_name = t.get("shape_name", "").lower()
                if slide_entry["title"] is None and ("title" in shape_name or len(slide_entry["text"]) == 0):
                    # Check if it looks like a title (short, no bullet points)
                    if len(text_content) < 200 and "\n" not in text_content:
                        slide_entry["title"] = text_content
                        continue
                slide_entry["text"].append(text_content)
            else:
                slide_entry["text"].append(t)
                
        if verbose:
            print(f"[DEBUG]   - Found {len(text_shapes)} text shapes")

        # 2) speaker notes
        if slide.has_notes_slide:
            try:
                notes = slide.notes_slide.notes_text_frame.text
                if notes and notes.strip():
                    if compact:
                        slide_entry["text"].append(f"[Notes]: {notes.strip()}")
                    else:
                        slide_entry["text"].append({"type": "notes", "text": notes})
            except Exception:
                pass

        # 3) tables (optional)
        if include_tables:
            table_cells = list(iter_table_cells(slide))
            if table_cells:
                if compact:
                    # Combine table cells into readable format
                    table_text = " | ".join(c.get("text", "") for c in table_cells if c.get("text"))
                    if table_text:
                        slide_entry["text"].append(f"[Table]: {table_text}")
                else:
                    for c in table_cells:
                        slide_entry["text"].append(c)
                        
            if verbose:
                print(f"[DEBUG]   - Found {len(table_cells)} table cells")

        # 4) images -> send to Document Intelligence for OCR
        SUPPORTED_IMAGE_FORMATS = {'jpg', 'jpeg', 'png', 'bmp', 'tiff', 'tif', 'heif', 'heic', 'pdf'}
        
        images = list(iter_images(slide))
        if verbose:
            print(f"[DEBUG]   - Found {len(images)} images to OCR")
        
        for idx, img in enumerate(images, start=1):
            ext = (img["ext"] or "png").lower().replace(".", "")
            name = f"slide{i}_img{idx}_{uuid.uuid4().hex[:8]}.{ext}"
            
            # Skip unsupported image formats (e.g., WMF, EMF)
            if ext not in SUPPORTED_IMAGE_FORMATS:
                if verbose:
                    print(f"[DEBUG]   - Skipping image {idx}/{len(images)}: {name} (unsupported format: {ext})")
                unsupported_images.append({
                    "slide": i,
                    "image_name": name,
                    "format": ext,
                    "reason": f"Unsupported format: {ext}"
                })
                if compact:
                    slide_entry["images"].append({
                        "name": name,
                        "ocr_text": "",
                        "skipped": True
                    })
                else:
                    slide_entry["images"].append({
                        "name": name,
                        "ocr": {"text": "", "skipped": True, "reason": f"Unsupported format: {ext}"},
                        "shape_name": img.get("shape_name"),
                    })
                continue
                
            ctype = f"image/{'jpeg' if ext in ('jpg', 'jpeg') else 'png' if ext=='png' else ext}"

            if verbose:
                print(f"[DEBUG]   - OCR image {idx}/{len(images)}: {name} ({len(img['bytes'])} bytes, {ctype})")
            
            try:
                # Use Document Intelligence for OCR
                ocr = analyze_image_bytes(di, img["bytes"], ctype)
                ocr_text = ocr.get("text", "")
                
                if verbose:
                    text_preview = ocr_text[:50].replace("\n", " ")
                    print(f"[DEBUG]     -> OCR complete: '{text_preview}...'")
            except Exception as e:
                if verbose:
                    print(f"[DEBUG]     -> OCR failed: {e}")
                ocr_text = ""
                ocr = {"text": "", "error": str(e)}

            if compact:
                slide_entry["images"].append({
                    "name": name,
                    "ocr_text": ocr_text
                })
            else:
                slide_entry["images"].append({
                    "name": name,
                    "ocr": ocr,
                    "shape_name": img.get("shape_name"),
                })

        slides_out.append(slide_entry)
        if verbose:
            print(f"[DEBUG] Slide {i} complete")

    if verbose:
        print(f"[DEBUG] All slides processed")
        if unsupported_images:
            print(f"[DEBUG] {len(unsupported_images)} unsupported images found")
    
    return {"slides": slides_out, "unsupported_images": unsupported_images}
