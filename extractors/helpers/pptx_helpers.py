"""PPTX shape extraction helpers."""
from __future__ import annotations
from typing import Dict, Any, Iterable

from pptx.enum.shapes import MSO_SHAPE_TYPE


def iter_text_shapes(slide) -> Iterable[Dict[str, Any]]:
    """Iterate over text shapes in a slide and yield text content."""
    for shp in slide.shapes:
        if getattr(shp, "has_text_frame", False):
            txt = shp.text_frame.text or ""
            if txt.strip():
                yield {
                    "type": "textbox",
                    "text": txt,
                    "shape_name": getattr(shp, "name", None),
                }


def iter_table_cells(slide) -> Iterable[Dict[str, Any]]:
    """Iterate over table cells in a slide and yield cell content."""
    for shp in slide.shapes:
        if shp.shape_type == MSO_SHAPE_TYPE.TABLE:
            t = shp.table
            for r_idx, row in enumerate(t.rows, start=1):
                for c_idx, cell in enumerate(row.cells, start=1):
                    text = cell.text or ""
                    if text.strip():
                        yield {
                            "type": "table_cell",
                            "row": r_idx,
                            "col": c_idx,
                            "text": text,
                            "shape_name": getattr(shp, "name", None),
                        }


def iter_images(slide) -> Iterable[Dict[str, Any]]:
    """Iterate over images in a slide and yield image data."""
    for shp in slide.shapes:
        if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
            img = shp.image
            ext = img.ext or "bin"
            yield {
                "type": "image",
                "ext": ext,
                "blob": img.blob,  # raw bytes (was "bytes", fixed to "blob" for multimodal_extract.py)
                "shape_name": getattr(shp, "name", None),
            }
