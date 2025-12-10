"""
Multimodal PPTX Extraction Pipeline

Combines ALL THREE sources for optimal extraction:
1. Native text extraction from PPTX (cleanest for text boxes)
2. Document Intelligence OCR on embedded images (catches text in screenshots)
3. LLM vision (GPT-4.1) to validate, structure, and reconcile all sources

Output: Clean, simplified JSON per slide with the most accurate text content.
"""
import os
import json
from pathlib import Path
from typing import Dict, Any, List, Optional, Union
from dataclasses import dataclass

from .config import DIConfig
from .llm_helpers import LLMConfig, analyze_slide_multimodal
from .slide_renderer import (
    render_slides_with_libreoffice, 
    render_slides_with_powerpoint,
    check_rendering_available
)
from .pptx_helpers import iter_text_shapes, iter_table_cells, iter_images
from .di_helpers import analyze_image_bytes, normalize_di_result
from .cache_storage import (
    CacheStorage,
    get_cache_storage,
    reset_cache_storage,
    compute_file_hash,
    get_cache_key
)

from pptx import Presentation
import sys


def _load_from_cache(
    pptx_path: str,
    model: str,
    use_di: bool,
    allow_local_cache: bool = False,
    verbose: bool = False
) -> Optional[Dict[str, Any]]:
    """
    Check if we have cached extraction results for this file.
    
    Returns cached results if found and valid, None otherwise.
    """
    if not os.path.exists(pptx_path):
        return None
    
    storage = get_cache_storage(allow_local=allow_local_cache, verbose=False)
    if not storage.is_available:
        return None
    
    file_hash = compute_file_hash(pptx_path)
    cache_key = get_cache_key(file_hash, prefix="pptx", model=model, di=str(use_di))
    
    cached = storage.get(cache_key)
    if cached is not None:
        # Verify cache metadata matches
        if cached.get("_cache_meta", {}).get("file_hash") == file_hash:
            if verbose:
                print(f"[Pipeline] Cache hit! Loading from cache")
            return cached
    
    return None


def _save_to_cache(
    results: Dict[str, Any],
    pptx_path: str,
    model: str,
    use_di: bool,
    allow_local_cache: bool = False,
    verbose: bool = False
) -> None:
    """Save extraction results to cache."""
    storage = get_cache_storage(allow_local=allow_local_cache, verbose=False)
    if not storage.is_available:
        return
    
    file_hash = compute_file_hash(pptx_path)
    cache_key = get_cache_key(file_hash, prefix="pptx", model=model, di=str(use_di))
    
    # Add cache metadata
    results["_cache_meta"] = {
        "file_hash": file_hash,
        "source_file": os.path.basename(pptx_path),
        "model": model,
        "use_di": use_di
    }
    
    storage.set(cache_key, results)
    
    if verbose:
        storage_name = type(storage).__name__
        print(f"[Pipeline] Results cached ({storage_name}, key: {cache_key[:20]}...)")


@dataclass
class MultimodalConfig:
    """Configuration for the multimodal extraction pipeline."""
    llm: LLMConfig                  # Primary LLM for text extraction from images
    di: Optional[DIConfig] = None   # Document Intelligence for embedded image OCR
    render_dpi: int = 150           # Slide render quality
    cache_rendered_slides: bool = True
    use_di_for_images: bool = True  # Whether to OCR embedded images with DI
    model_type: str = "gpt-4.1"     # Model type: "gpt-4.1" or "gpt-5.1"
    fallback_llm: Optional[LLMConfig] = None  # Fallback LLM if primary returns empty


def extract_native_text(slide, include_tables: bool = True) -> str:
    """Extract native text from a slide (no OCR)."""
    text_parts = []
    
    # Native text boxes
    for t in iter_text_shapes(slide):
        if t.get("text"):
            text_parts.append(t["text"])
    
    # Speaker notes
    if slide.has_notes_slide:
        try:
            notes = slide.notes_slide.notes_text_frame.text
            if notes and notes.strip():
                text_parts.append(f"[Speaker Notes]: {notes}")
        except Exception:
            pass
    
    # Tables
    if include_tables:
        table_text = []
        for c in iter_table_cells(slide):
            if c.get("text"):
                table_text.append(c["text"])
        if table_text:
            text_parts.append(f"[Table]: {' | '.join(table_text)}")
    
    return "\n\n".join(text_parts)


# Supported image formats for Document Intelligence
SUPPORTED_IMAGE_FORMATS = {"jpg", "jpeg", "png", "bmp", "tiff", "tif", "heif", "heic", "pdf"}


def extract_di_ocr_from_images(slide, di_config: DIConfig, slide_index: int) -> str:
    """Extract OCR text from embedded images using Document Intelligence."""
    ocr_parts = []
    
    for img in iter_images(slide):
        # Check if format is supported
        ext = img.get("ext", "").lower().lstrip(".")
        if ext not in SUPPORTED_IMAGE_FORMATS:
            continue
            
        img_bytes = img.get("blob")
        if not img_bytes:
            continue
        
        try:
            result = analyze_image_bytes(di_config, img_bytes, ext)
            normalized = normalize_di_result(result, compact=True)
            text = normalized.get("text", "").strip()
            if text:
                ocr_parts.append(f"[Image OCR]: {text}")
        except Exception:
            # Skip failed OCR silently
            pass
    
    return "\n\n".join(ocr_parts)


def multimodal_extract(
    pptx_path: str,
    config: MultimodalConfig,
    output_path: Optional[str] = None,
    verbose: bool = True
) -> Dict[str, Any]:
    """
    Extract text from a PPTX using multimodal LLM analysis.
    
    Pipeline:
    1. Render all slides to images
    2. Extract native text from each slide (for reference)
    3. Send image + text to GPT-4.1 for accurate text extraction
    4. Return simplified JSON
    
    Args:
        pptx_path: Path to the PPTX file
        config: Multimodal configuration
        output_path: Optional path to save output JSON
        verbose: Print progress
        
    Returns:
        Simplified JSON with slides containing index, title, and text content
    """
    import tempfile
    
    # Check rendering availability
    can_render, msg = check_rendering_available()
    if not can_render:
        raise RuntimeError(f"Slide rendering not available: {msg}")
    
    if verbose:
        print(f"[Pipeline] Opening: {pptx_path}")
    
    prs = Presentation(pptx_path)
    total_slides = len(prs.slides)
    
    if verbose:
        print(f"[Pipeline] Found {total_slides} slides")
    
    # Step 1: Render slides to images
    if verbose:
        print(f"[Pipeline] Rendering slides to images...")
    
    if config.cache_rendered_slides:
        cache_dir = os.path.join(os.path.dirname(pptx_path), ".slide_cache")
        os.makedirs(cache_dir, exist_ok=True)
    else:
        cache_dir = tempfile.mkdtemp()
    
    # Use PowerPoint on Windows if available, else LibreOffice
    try:
        if sys.platform == "win32":
            try:
                import comtypes.client
                slide_images = render_slides_with_powerpoint(pptx_path, cache_dir, config.render_dpi)
                if verbose:
                    print(f"[Pipeline] Using PowerPoint COM for rendering")
            except Exception as e:
                if verbose:
                    print(f"[Pipeline] PowerPoint not available, trying LibreOffice: {e}")
                slide_images = render_slides_with_libreoffice(pptx_path, cache_dir, config.render_dpi)
        else:
            slide_images = render_slides_with_libreoffice(pptx_path, cache_dir, config.render_dpi)
    except Exception as e:
        raise RuntimeError(f"Failed to render slides: {e}")
    
    if verbose:
        print(f"[Pipeline] Rendered {len(slide_images)} slide images")
    
    # Check if DI is configured for embedded image OCR
    use_di = config.use_di_for_images and config.di is not None
    if verbose:
        if use_di:
            print(f"[Pipeline] Document Intelligence enabled for embedded image OCR")
        else:
            print(f"[Pipeline] Document Intelligence not configured (using native text only)")
    
    # Step 2 & 3: Process each slide with all three sources
    results = []
    
    for i, slide in enumerate(prs.slides, start=1):
        if verbose:
            print(f"[Pipeline] Processing slide {i}/{total_slides}...")
        
        # Source 1: Native text from PPTX (cleanest for text boxes)
        native_text = extract_native_text(slide)
        
        # Source 2: DI OCR from embedded images (catches text in screenshots)
        di_ocr_text = ""
        if use_di:
            try:
                di_ocr_text = extract_di_ocr_from_images(slide, config.di, i)
                if di_ocr_text and verbose:
                    print(f"[Pipeline]   DI OCR extracted from embedded images")
            except Exception as e:
                if verbose:
                    print(f"[Pipeline]   DI OCR failed: {e}")
        
        # Combine native text + DI OCR for LLM context
        combined_extracted_text = native_text
        if di_ocr_text:
            combined_extracted_text = f"{native_text}\n\n{di_ocr_text}"
        
        # Source 3: LLM vision on rendered slide image
        if i <= len(slide_images):
            with open(slide_images[i - 1], "rb") as f:
                slide_image_bytes = f.read()
        else:
            if verbose:
                print(f"[Pipeline]   Warning: No image for slide {i}")
            slide_image_bytes = None
        
        # Send ALL sources to LLM for final reconciliation
        if slide_image_bytes:
            try:
                # GPT-5.1 requires max_completion_tokens instead of max_tokens
                use_max_completion_tokens = config.model_type == "gpt-5.1"
                
                llm_text = analyze_slide_multimodal(
                    config.llm,
                    slide_image_bytes,
                    combined_extracted_text,  # Now includes both native + DI OCR
                    use_max_completion_tokens=use_max_completion_tokens
                )
                
                # Fallback to secondary LLM if primary returns empty
                if not llm_text and config.fallback_llm:
                    if verbose:
                        print(f"[Pipeline]   Primary LLM returned empty, trying fallback...")
                    llm_text = analyze_slide_multimodal(
                        config.fallback_llm,
                        slide_image_bytes,
                        combined_extracted_text,
                        use_max_completion_tokens=False  # GPT-4.1 uses max_tokens
                    )
                
                if verbose:
                    preview = llm_text[:60].replace("\n", " ") if llm_text else "(empty)"
                    print(f"[Pipeline]   -> Extracted: {preview}...")
                    
            except Exception as e:
                if verbose:
                    print(f"[Pipeline]   Error analyzing slide {i}: {e}")
                llm_text = combined_extracted_text  # Fallback to combined text
        else:
            llm_text = combined_extracted_text
        
        results.append({
            "index": i,
            "source_file": os.path.basename(pptx_path),
            "source_index": i,
            "text": llm_text
        })
    
    output = {
        "source_files": [os.path.basename(pptx_path)],
        "total_slides": total_slides,
        "slides": results
    }
    
    # Save if requested
    if output_path:
        with open(output_path, "w", encoding="utf-8") as f:
            json.dump(output, f, indent=2, ensure_ascii=False)
        if verbose:
            print(f"[Pipeline] Output saved to: {output_path}")
    
    return output


def quick_extract(
    pptx_path: str,
    output_path: Optional[str] = None,
    verbose: bool = True,
    use_di: bool = True,
    model: str = "gpt-4.1",
    use_cache: bool = True,
    allow_local_cache: bool = False
) -> Dict[str, Any]:
    """
    Convenience function that loads config from environment variables.
    
    Combines ALL THREE sources for optimal extraction:
    1. Native text from PPTX (cleanest for text boxes)
    2. Document Intelligence OCR on embedded images (catches text in screenshots)
    3. LLM vision to validate and reconcile all sources
    
    Args:
        pptx_path: Path to the PPTX file
        output_path: Optional path to save output JSON
        verbose: Print progress
        use_di: Whether to use Document Intelligence for embedded image OCR
        model: Which model to use - "gpt-4.1" or "gpt-5.1"
        use_cache: Whether to use cached results if available (default True)
        allow_local_cache: Allow local filesystem cache (for development only)
    
    Required env vars for GPT-4.1:
    - AZURE_AI_ENDPOINT, AZURE_AI_API_KEY, GPT_4_1_DEPLOYMENT
    
    Required env vars for GPT-5.1:
    - AZURE_AI_GPT5_ENDPOINT, AZURE_AI_GPT5_API_KEY, GPT_5_1_DEPLOYMENT
    
    Optional env vars (for DI OCR):
    - AZURE_DI_ENDPOINT, AZURE_DI_KEY
    
    Optional env vars (for Azure Blob cache):
    - AZURE_STORAGE_CONNECTION_STRING or AZURE_STORAGE_ACCOUNT_NAME
    """
    from dotenv import load_dotenv
    load_dotenv()
    
    # Check cache first
    if use_cache:
        cached = _load_from_cache(pptx_path, model, use_di, allow_local_cache, verbose)
        if cached is not None:
            # Still save to output_path if requested
            if output_path:
                with open(output_path, "w", encoding="utf-8") as f:
                    json.dump(cached, f, indent=2, ensure_ascii=False)
                if verbose:
                    print(f"[Pipeline] Cached output copied to: {output_path}")
            return cached
    
    # LLM config based on model selection
    fallback_llm = None
    if model == "gpt-5.1":
        llm_config = LLMConfig(
            endpoint=os.environ["AZURE_AI_GPT5_ENDPOINT"],
            api_key=os.environ["AZURE_AI_GPT5_API_KEY"],
            deployment=os.environ["GPT_5_1_DEPLOYMENT"]
        )
        # Set up GPT-4.1 as fallback for GPT-5.1
        if os.getenv("AZURE_AI_ENDPOINT") and os.getenv("AZURE_AI_API_KEY") and os.getenv("GPT_4_1_DEPLOYMENT"):
            fallback_llm = LLMConfig(
                endpoint=os.environ["AZURE_AI_ENDPOINT"],
                api_key=os.environ["AZURE_AI_API_KEY"],
                deployment=os.environ["GPT_4_1_DEPLOYMENT"]
            )
            if verbose:
                print(f"[Pipeline] Using GPT-5.1 model with GPT-4.1 fallback")
        else:
            if verbose:
                print(f"[Pipeline] Using GPT-5.1 model (no fallback configured)")
    else:
        llm_config = LLMConfig(
            endpoint=os.environ["AZURE_AI_ENDPOINT"],
            api_key=os.environ["AZURE_AI_API_KEY"],
            deployment=os.environ["GPT_4_1_DEPLOYMENT"]
        )
        if verbose:
            print(f"[Pipeline] Using GPT-4.1 model")
    
    # DI config (optional - for embedded image OCR)
    di_config = None
    if use_di and os.getenv("AZURE_DI_ENDPOINT") and os.getenv("AZURE_DI_KEY"):
        di_config = DIConfig(
            endpoint=os.environ["AZURE_DI_ENDPOINT"],
            key=os.environ["AZURE_DI_KEY"]
        )
    
    config = MultimodalConfig(
        llm=llm_config,
        di=di_config,
        use_di_for_images=use_di and di_config is not None,
        model_type=model,
        fallback_llm=fallback_llm
    )
    
    results = multimodal_extract(pptx_path, config, output_path, verbose)
    
    # Save to cache for future runs
    if use_cache:
        _save_to_cache(results, pptx_path, model, use_di, allow_local_cache, verbose)
    
    return results


def quick_extract_gpt5(
    pptx_path: str,
    output_path: Optional[str] = None,
    verbose: bool = True,
    use_di: bool = True
) -> Dict[str, Any]:
    """
    Convenience function for GPT-5.1 extraction.
    
    Same as quick_extract but uses GPT-5.1 by default.
    
    Required env vars:
    - AZURE_AI_GPT5_ENDPOINT, AZURE_AI_GPT5_API_KEY, GPT_5_1_DEPLOYMENT
    """
    return quick_extract(pptx_path, output_path, verbose, use_di, model="gpt-5.1")


def quick_extract_multi(
    pptx_paths: List[str],
    output_path: Optional[str] = None,
    verbose: bool = True,
    use_di: bool = True,
    model: str = "gpt-4.1",
    use_cache: bool = True,
    allow_local_cache: bool = False
) -> Dict[str, Any]:
    """
    Extract text from MULTIPLE PPTX files using multimodal LLM analysis.
    
    Processes each PPTX in order, combining results with continuous slide indexing.
    Each slide includes source_file and source_index for traceability.
    Uses per-file caching so unchanged files won't be re-extracted.
    
    Args:
        pptx_paths: List of paths to PPTX files (processed in order)
        output_path: Optional path to save combined output JSON
        verbose: Print progress
        use_di: Whether to use Document Intelligence for embedded image OCR
        model: Which model to use - "gpt-4.1" or "gpt-5.1"
        use_cache: Whether to use cached results if available (default True)
        allow_local_cache: Allow local filesystem cache (for development only)
    
    Returns:
        Combined JSON with all slides from all files:
        {
            "source_files": ["file1.pptx", "file2.pptx"],
            "total_slides": 120,
            "slides": [
                {"index": 1, "source_file": "file1.pptx", "source_index": 1, "text": "..."},
                {"index": 70, "source_file": "file2.pptx", "source_index": 1, "text": "..."}
            ]
        }
    """
    if not pptx_paths:
        raise ValueError("At least one PPTX path is required")
    
    if len(pptx_paths) == 1:
        # Single file - just use quick_extract
        return quick_extract(pptx_paths[0], output_path, verbose, use_di, model, use_cache, allow_local_cache)
    
    if verbose:
        print(f"[Multi-PPTX] Processing {len(pptx_paths)} files...")
    
    all_slides = []
    source_files = []
    global_index = 0
    
    for file_num, pptx_path in enumerate(pptx_paths, start=1):
        if verbose:
            print(f"\n[Multi-PPTX] File {file_num}/{len(pptx_paths)}: {os.path.basename(pptx_path)}")
        
        # Extract from this file (without saving) - uses per-file caching
        result = quick_extract(pptx_path, None, verbose, use_di, model, use_cache, allow_local_cache)
        
        source_files.append(os.path.basename(pptx_path))
        
        # Re-index slides with global continuous numbering
        for slide in result["slides"]:
            global_index += 1
            slide["index"] = global_index  # Global continuous index
            # source_file and source_index are already set by multimodal_extract
            all_slides.append(slide)
    
    combined_output = {
        "source_files": source_files,
        "total_slides": global_index,
        "slides": all_slides
    }
    
    # Save if requested
    if output_path:
        with open(output_path, "w", encoding="utf-8") as f:
            json.dump(combined_output, f, indent=2, ensure_ascii=False)
        if verbose:
            print(f"\n[Multi-PPTX] Combined output saved to: {output_path}")
    
    if verbose:
        print(f"[Multi-PPTX] Total: {global_index} slides from {len(source_files)} files")
    
    return combined_output


# ==============================================================================
# CLI Interface
# ==============================================================================

def main():
    """CLI entry point for multimodal PPTX extraction."""
    import argparse
    
    parser = argparse.ArgumentParser(
        description="Extract text from PPTX files using multimodal LLM analysis (native + DI OCR + GPT vision)",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  # Single file extraction
  python -m extractors.helpers.multimodal_extract evidence.pptx -o output.json
  
  # Multiple files (combined output)
  python -m extractors.helpers.multimodal_extract file1.pptx file2.pptx file3.pptx -o combined.json
  
  # Use GPT-5.1 instead of GPT-4.1
  python -m extractors.helpers.multimodal_extract evidence.pptx -o output.json --model gpt-5.1
  
  # Skip Document Intelligence OCR
  python -m extractors.helpers.multimodal_extract evidence.pptx -o output.json --no-di

Required environment variables:
  AZURE_AI_ENDPOINT, AZURE_AI_API_KEY, GPT_4_1_DEPLOYMENT (for GPT-4.1)
  AZURE_AI_GPT5_ENDPOINT, AZURE_AI_GPT5_API_KEY, GPT_5_1_DEPLOYMENT (for GPT-5.1)
  AZURE_DI_ENDPOINT, AZURE_DI_KEY (optional, for embedded image OCR)
"""
    )
    
    parser.add_argument(
        "pptx_files",
        nargs="+",
        help="One or more PPTX files to process"
    )
    parser.add_argument(
        "-o", "--output",
        help="Output JSON file path (default: prints to stdout)"
    )
    parser.add_argument(
        "--model",
        choices=["gpt-4.1", "gpt-5.1"],
        default="gpt-4.1",
        help="Which model to use (default: gpt-4.1)"
    )
    parser.add_argument(
        "--no-di",
        action="store_true",
        help="Skip Document Intelligence OCR for embedded images"
    )
    parser.add_argument(
        "-q", "--quiet",
        action="store_true",
        help="Suppress progress output"
    )
    
    args = parser.parse_args()
    
    # Validate input files exist
    from pathlib import Path
    for pptx_path in args.pptx_files:
        if not Path(pptx_path).exists():
            print(f"Error: File not found: {pptx_path}", file=sys.stderr)
            sys.exit(1)
        if not pptx_path.lower().endswith(('.pptx', '.ppt')):
            print(f"Warning: {pptx_path} may not be a PowerPoint file", file=sys.stderr)
    
    verbose = not args.quiet
    use_di = not args.no_di
    
    try:
        if len(args.pptx_files) == 1:
            result = quick_extract(
                args.pptx_files[0],
                args.output,
                verbose,
                use_di,
                args.model
            )
        else:
            result = quick_extract_multi(
                args.pptx_files,
                args.output,
                verbose,
                use_di,
                args.model
            )
        
        # If no output file, print to stdout
        if not args.output:
            print(json.dumps(result, indent=2, ensure_ascii=False))
            
    except KeyError as e:
        print(f"Error: Missing environment variable: {e}", file=sys.stderr)
        sys.exit(1)
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
